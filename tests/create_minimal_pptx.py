from __future__ import annotations

from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "examples" / "demo.pptx"


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(5.8), Inches(0.7))
    run = text_box.text_frame.paragraphs[0].add_run()
    run.text = "我们需要提升内容质量"
    run.font.size = Pt(28)

    image_stream = BytesIO()
    Image.new("RGB", (100, 80), (42, 132, 214)).save(image_stream, format="PNG")
    image_stream.seek(0)
    slide.shapes.add_picture(image_stream, Inches(1.1), Inches(1.5), width=Inches(2.2))

    presentation.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
